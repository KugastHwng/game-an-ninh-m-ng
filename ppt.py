from pptx import Presentation
from pptx.util import Inches, Pt
import json
import requests

# Link raw JSON trên GitHub (thay bằng link của bạn)
json_url = "https://raw.githubusercontent.com/KugastHwng/nh-p/refs/heads/main/scores.json"

response = requests.get(json_url)
if response.status_code == 200:
    players = response.json()
else:
    print("Không lấy được dữ liệu từ GitHub. Dùng file local 'scores_100.json'")
    with open('scores_100.json', 'r') as f:
        players = json.load(f)

prs = Presentation()

# Slide tiêu đề
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
slide_title.shapes.title.text = "Bảng Kết Quả Trò Chơi"
slide_title.placeholders[1].text = f"Số người chơi: {len(players)}"

# Hiển thị mỗi nhóm 10 người / slide
group_size = 10
for i in range(0, len(players), group_size):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    top = Inches(1)
    left = Inches(1)
    for j, player in enumerate(players[i:i+group_size]):
        txBox = slide.shapes.add_textbox(left, top + Inches(j*0.5), Inches(6), Inches(0.5))
        tf = txBox.text_frame
        tf.text = f"{player['name']} - Số câu đúng: {player['score']}"
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)

prs.save('KetQua_100.pptx')
print("Đã tạo PowerPoint: KetQua_100.pptx")
